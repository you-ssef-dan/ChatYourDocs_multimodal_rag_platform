# text_ingestion.py
import os
import fitz  # PyMuPDF for PDF image extraction
from tqdm import tqdm
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document as DocxDocument
from openpyxl import load_workbook

from langchain_experimental.text_splitter import SemanticChunker
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_unstructured import UnstructuredLoader
from langchain.document_loaders import PyPDFLoader
from langchain_community.document_loaders import (
    UnstructuredExcelLoader,
    UnstructuredCSVLoader,
    UnstructuredWordDocumentLoader,
)
from langchain.schema import Document

from vectordb import vector_db
from images_ingestion import ingest_images


print("🔧 Initializing text ingestion...")


# ---------- Image Extraction Helpers ----------

def extract_images_from_pdf(path, images_folder="documents"):
    os.makedirs(images_folder, exist_ok=True)
    doc = fitz.open(path)
    image_paths = []
    base_name = os.path.splitext(os.path.basename(path))[0]
    count = 0

    for page_num, page in enumerate(doc):
        for img_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            ext = base_image["ext"]
            count += 1
            filename = f"{base_name}_page{page_num+1}_img{count}.{ext}"
            filepath = os.path.join(images_folder, filename)
            with open(filepath, "wb") as f:
                f.write(image_bytes)
            image_paths.append(filepath)

    return image_paths


def extract_images_from_docx(path, images_folder="documents"):
    doc = DocxDocument(path)
    image_paths = []
    count = 0
    base_name = os.path.splitext(os.path.basename(path))[0]

    for rel in doc.part._rels:
        rel_target = doc.part._rels[rel]
        if "image" in rel_target.target_ref:
            count += 1
            image_bytes = rel_target.target_part.blob
            filename = f"{base_name}_image_{count}.png"
            filepath = os.path.join(images_folder, filename)
            with open(filepath, "wb") as f:
                f.write(image_bytes)
            image_paths.append(filepath)

    return image_paths


def extract_images_from_pptx(path, images_folder="documents"):
    os.makedirs(images_folder, exist_ok=True)
    prs = Presentation(path)
    image_paths = []
    count = 0
    base_name = os.path.splitext(os.path.basename(path))[0]

    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                count += 1
                image = shape.image
                ext = image.ext
                filename = f"{base_name}_slide{slide_num}_img{count}.{ext}"
                filepath = os.path.join(images_folder, filename)
                with open(filepath, "wb") as f:
                    f.write(image.blob)
                image_paths.append(filepath)

    return image_paths


def extract_images_from_excel(path, images_folder="documents"):
    os.makedirs(images_folder, exist_ok=True)
    wb = load_workbook(path, data_only=True)
    image_paths = []
    count = 0
    base_name = os.path.splitext(os.path.basename(path))[0]

    for sheet in wb.worksheets:
        for image in getattr(sheet, "_images", []):  # openpyxl stores images here
            count += 1
            filename = f"{base_name}_sheet{sheet.title}_img{count}.png"
            filepath = os.path.join(images_folder, filename)
            with open(filepath, "wb") as f:
                f.write(image._data())
            image_paths.append(filepath)

    return image_paths


# ---------- Text Loaders ----------

def load_pptx(path):
    prs = Presentation(path)
    docs = []

    for i, slide in enumerate(prs.slides):
        slide_texts = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
            elif shape.has_table:
                table = shape.table
                rows = []
                for r in range(len(table.rows)):
                    cells = [table.cell(r, c).text.strip() for c in range(len(table.columns))]
                    rows.append("\t".join(cells))
                slide_texts.append("\n".join(rows))

        content = "\n\n".join(slide_texts).strip()
        if content:
            docs.append(Document(page_content=content, metadata={"source": path, "slide": i + 1}))

    return docs


FILE_LOADERS = {
    ".pdf": lambda p: PyPDFLoader(p),
    ".txt": UnstructuredLoader,
    ".docx": UnstructuredWordDocumentLoader,
    ".doc": UnstructuredWordDocumentLoader,
    ".xlsx": UnstructuredExcelLoader,
    ".xls": UnstructuredExcelLoader,
    ".csv": UnstructuredCSVLoader,
    ".pptx": load_pptx,
    ".ppt": load_pptx,
    ".html": UnstructuredLoader,
    ".htm": UnstructuredLoader,
    ".rtf": UnstructuredLoader,
    ".odt": UnstructuredLoader,
    ".msg": UnstructuredLoader,
}


# ---------- Document Loader with Image Extraction ----------

def load_documents(file_paths):
    docs = []
    extracted_images = []

    for path in tqdm(file_paths, desc="📂 Loading documents"):
        ext = os.path.splitext(path)[1].lower()

        try:
            if ext in [".docx", ".doc"]:
                loader = UnstructuredWordDocumentLoader(path)
                loaded_docs = loader.load()
                docs.extend(loaded_docs)
                extracted_images.extend(extract_images_from_docx(path))

            elif ext == ".pdf":
                loader = PyPDFLoader(path)
                loaded_docs = loader.load()
                docs.extend(loaded_docs)
                extracted_images.extend(extract_images_from_pdf(path))

            elif ext in [".pptx", ".ppt"]:
                loaded_docs = load_pptx(path)
                docs.extend(loaded_docs)
                extracted_images.extend(extract_images_from_pptx(path))

            elif ext in [".xlsx", ".xls"]:
                loader = UnstructuredExcelLoader(path)
                loaded_docs = loader.load()
                docs.extend(loaded_docs)
                extracted_images.extend(extract_images_from_excel(path))

            else:
                loader_fn = FILE_LOADERS.get(ext)
                if loader_fn:
                    loaded_docs = loader_fn(path)
                    if callable(getattr(loaded_docs, "load", None)):
                        loaded_docs = loaded_docs.load()
                    docs.extend(loaded_docs)

        except Exception as e:
            print(f"❌ Failed to load {path}: {e}")

    print(f"📊 Total documents loaded: {len(docs)} chunks")
    print(f"📊 Total extracted images: {len(extracted_images)}")

    return docs, extracted_images


# ---------- Main Ingestion Function ----------

def ingest_texts(collection_name="documents", file_paths=None):
    if not file_paths:
        print("⚠️ No documents to ingest!")
        return

    docs, extracted_images = load_documents(file_paths)

    if not docs:
        print("⚠️ No documents were successfully loaded.")
        return

    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-mpnet-base-v2")
    splitter = SemanticChunker(embeddings=embeddings, breakpoint_threshold_type="percentile")
    chunks = splitter.split_documents(docs)

    vector_db.store_documents(collection_name=collection_name, documents=chunks)
    print(f"✅ {len(chunks)} text chunks added to '{collection_name}'.")

    if extracted_images:
        print(f"🖼️ Ingesting {len(extracted_images)} images…")
        ingest_images(collection_name="images", file_paths=extracted_images)
